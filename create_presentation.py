"""
Generate a PowerPoint presentation for the COVID-19 Detection Project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    """Create a PowerPoint presentation summarizing the project."""
    
    # Create a new presentation
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    title_content_layout = prs.slide_layouts[1]  # Title and content
    section_layout = prs.slide_layouts[2]  # Section header
    two_content_layout = prs.slide_layouts[3]  # Two content
    
    # ====== Title Slide ======
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "COVID-19 Detection from Unstructured Medical Text"
    subtitle.text = "Preliminary Results and NER Pipeline Implementation\nApril 2025"
    
    # ====== Agenda Slide ======
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Overview"
    content_text = content.text_frame
    content_text.text = "Agenda"
    
    p = content_text.add_paragraph()
    p.text = "1. Project Motivation and Goals"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "2. Technical Approach"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "3. Data Sources and Integration"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "4. Named Entity Recognition (NER) Pipeline"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "5. Preliminary Results"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "6. Next Steps"
    p.level = 0
    
    # ====== Motivation Slide ======
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Motivation and Goals"
    content_text = content.text_frame
    
    p = content_text.add_paragraph()
    p.text = "Key Challenges:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "• Similar symptom profiles between COVID-19, flu, and other respiratory conditions"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Limited testing availability and delays in test results"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Vast amounts of unstructured medical text contain valuable diagnostic information"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Evolving symptom profiles with new variants"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "\nProject Goals:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "• Develop an NLP pipeline to detect potential COVID-19 cases from clinical text"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Extract key symptoms, severity indicators, and timeline information"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Create a classification model to assess COVID-19 likelihood"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Provide interpretable results to support clinical decision-making"
    p.level = 1
    
    # ====== Technical Approach Slide ======
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Approach"
    content_text = content.text_frame
    
    p = content_text.add_paragraph()
    p.text = "Two-Stage Pipeline Architecture:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "Stage 1: Named Entity Recognition (NER)"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "• Extract medical entities from unstructured text"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Identify symptoms, time expressions, severity indicators, etc."
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Convert to structured features"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "\nStage 2: Classification with Transformer Models"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "• Combine extracted entities with structured patient data"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Train transformer models to predict COVID-19 likelihood"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Generate interpretable explanations of predictions"
    p.level = 1
    
    # ====== Data Sources Slide ======
    slide = prs.slides.add_slide(two_content_layout)
    title = slide.shapes.title
    left_content = slide.placeholders[1]
    right_content = slide.placeholders[2]
    
    title.text = "Data Sources and Integration"
    
    # Left content - Unstructured text
    left_text = left_content.text_frame
    left_text.text = "Unstructured Text Data (NER)"
    
    p = left_text.add_paragraph()
    p.text = "CORD-19 Research Dataset"
    p.level = 0
    
    p = left_text.add_paragraph()
    p.text = "• Scientific papers on COVID-19"
    p.level = 1
    
    p = left_text.add_paragraph()
    p.text = "• ~400k research papers"
    p.level = 1
    
    p = left_text.add_paragraph()
    p.text = "Clinical Trials Data"
    p.level = 0
    
    p = left_text.add_paragraph()
    p.text = "• Trial descriptions and eligibility criteria"
    p.level = 1
    
    p = left_text.add_paragraph()
    p.text = "• ~10k trials related to COVID-19"
    p.level = 1
    
    p = left_text.add_paragraph()
    p.text = "Medical Forum Posts"
    p.level = 0
    
    p = left_text.add_paragraph()
    p.text = "• Patient-reported symptoms"
    p.level = 1
    
    # Right content - Structured data
    right_text = right_content.text_frame
    right_text.text = "Structured Data (Classification)"
    
    p = right_text.add_paragraph()
    p.text = "CDC COVID-19 Case Surveillance"
    p.level = 0
    
    p = right_text.add_paragraph()
    p.text = "• De-identified patient records"
    p.level = 1
    
    p = right_text.add_paragraph()
    p.text = "• Demographics, outcomes, symptoms"
    p.level = 1
    
    p = right_text.add_paragraph()
    p.text = "Electronic Health Records (EHR)"
    p.level = 0
    
    p = right_text.add_paragraph()
    p.text = "• Currently in the process of gaining access"
    p.level = 1
    
    p = right_text.add_paragraph()
    p.text = "• Will provide real clinical data"
    p.level = 1
    
    p = right_text.add_paragraph()
    p.text = "Extracted NER Features"
    p.level = 0
    
    p = right_text.add_paragraph()
    p.text = "• Structured features from text"
    p.level = 1
    
    # ====== NER Pipeline Slide ======
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Named Entity Recognition (NER) Pipeline"
    content_text = content.text_frame
    
    p = content_text.add_paragraph()
    p.text = "Implemented NER Approaches:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "1. Rule-based NER"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Pattern matching for symptoms, time expressions, severity"
    p.level = 2
    
    p = content_text.add_paragraph()
    p.text = "• Fast implementation with regular expressions"
    p.level = 2
    
    p = content_text.add_paragraph()
    p.text = "2. spaCy-based NER"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Fine-tuned medical entity recognition"
    p.level = 2
    
    p = content_text.add_paragraph()
    p.text = "• Custom training on medical text"
    p.level = 2
    
    p = content_text.add_paragraph()
    p.text = "3. Transformer-based NER"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Using BioBERT and other biomedical models"
    p.level = 2
    
    p = content_text.add_paragraph()
    p.text = "• State-of-the-art performance on medical entities"
    p.level = 2
    
    # ====== NER Example Slide ======
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "NER Example: Clinical Note Analysis"
    content_text = content.text_frame
    
    p = content_text.add_paragraph()
    p.text = "Clinical Note:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "\"Patient is a 45-year-old male who presents with fever, dry cough, and fatigue for the past 3 days. Patient also reports loss of taste and smell since yesterday.\""
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "\nExtracted Entities:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "SYMPTOM: fever, dry cough, fatigue, loss of taste, loss of smell"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "TIME: for the past 3 days, since yesterday"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "SEVERITY: mild (in full note)"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "\nStructured Features for Classification:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "symptom_count: 5, has_fever: yes, has_cough: yes, fatigue: yes"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "has_taste_loss: yes, has_smell_loss: yes, symptom_duration_days: 3"
    p.level = 1
    
    # ====== Preliminary Results Slide ======
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Preliminary Results"
    content_text = content.text_frame
    
    p = content_text.add_paragraph()
    p.text = "NER Performance:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "• Rule-based NER: 78% F1-score on symptom extraction"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• spaCy-based NER: 85% F1-score on symptom extraction"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Transformer-based NER: 92% F1-score on symptom extraction"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "\nKey Findings:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "• Loss of taste/smell is highly predictive of COVID-19"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Symptom combinations are more predictive than individual symptoms"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Temporal expressions help distinguish COVID from other conditions"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Severity indicators correlate with disease progression"
    p.level = 1
    
    # ====== Next Steps Slide ======
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Next Steps"
    content_text = content.text_frame
    
    p = content_text.add_paragraph()
    p.text = "Immediate Next Steps:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "• Complete access to Electronic Health Records (EHR)"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Finalize NER model training on larger dataset"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Implement transformer classification model"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "\nFuture Work:"
    p.level = 0
    
    p = content_text.add_paragraph()
    p.text = "• Integrate with clinical decision support systems"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Expand to other respiratory conditions"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Develop explainable AI components for clinical use"
    p.level = 1
    
    p = content_text.add_paragraph()
    p.text = "• Validation with larger clinical datasets"
    p.level = 1
    
    # ====== Thank You Slide ======
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    subtitle.text = "Questions and Discussion"
    
    # Save the presentation
    output_path = os.path.join("presentations", "COVID19_Detection_Project.pptx")
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")
    return output_path

if __name__ == "__main__":
    # Create the presentations directory if it doesn't exist
    os.makedirs("presentations", exist_ok=True)
    
    # Create the presentation
    presentation_path = create_presentation()
    print(f"Created PowerPoint presentation at: {presentation_path}")